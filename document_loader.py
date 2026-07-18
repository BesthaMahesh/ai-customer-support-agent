import os
import io
import uuid
import pandas as pd
from pypdf import PdfReader
import docx2txt
from pptx import Presentation
from typing import List, Dict, Any
from utils import get_logger

logger = get_logger("DocumentLoader")

def parse_in_memory_file(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """Parses raw file bytes based on filename extension and returns a list of text page objects.
    
    Each page object is a dictionary:
      {"text": str, "page": int/str, "source": str}
    """
    ext = os.path.splitext(filename)[1].lower()
    logger.info(f"Parsing uploaded file: '{filename}' (Extension: {ext}, Size: {len(file_bytes)} bytes)")
    
    try:
        # 1. PDF Parser
        if ext == ".pdf":
            reader = PdfReader(io.BytesIO(file_bytes))
            pages = []
            for i, page in enumerate(reader.pages):
                try:
                    text = page.extract_text()
                    if text and text.strip():
                        pages.append({
                            "text": text,
                            "page": i + 1,
                            "source": filename
                        })
                except Exception as page_err:
                    logger.error(f"Error extracting text from PDF page {i+1}: {page_err}")
            
            if not pages:
                # Handle potential empty or scanned PDF exception gracefully
                logger.warning(f"PDF '{filename}' yielded no text. It might be scanned or empty.")
            return pages
            
        # 2. Text / Markdown / JSON Parser
        elif ext in [".txt", ".md", ".json"]:
            text = file_bytes.decode("utf-8", errors="ignore")
            if not text.strip():
                return []
            return [{
                "text": text,
                "page": 1,
                "source": filename
            }]
            
        # 3. CSV Table Parser
        elif ext == ".csv":
            df = pd.read_csv(io.BytesIO(file_bytes))
            if df.empty:
                return []
            # Format row representation: "Col1: Val1, Col2: Val2..."
            rows = []
            for idx, row in df.iterrows():
                row_str = ", ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                rows.append(f"Row {idx+1}: {row_str}")
            text = "\n".join(rows)
            return [{
                "text": text,
                "page": 1,
                "source": filename
            }]
            
        # 4. Excel spreadsheet sheet-by-sheet Parser
        elif ext == ".xlsx":
            xls = pd.ExcelFile(io.BytesIO(file_bytes))
            pages = []
            for sheet_name in xls.sheet_names:
                df = xls.parse(sheet_name)
                if df.empty:
                    continue
                rows = []
                for idx, row in df.iterrows():
                    row_str = ", ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                    rows.append(f"Row {idx+1}: {row_str}")
                text = f"Sheet: {sheet_name}\n" + "\n".join(rows)
                pages.append({
                    "text": text,
                    "page": sheet_name,
                    "source": filename
                })
            return pages
            
        # 5. Word Document (DOCX) Parser (using temp file to ensure library compatibility)
        elif ext == ".docx":
            temp_dir = "temp_uploads"
            os.makedirs(temp_dir, exist_ok=True)
            temp_path = os.path.join(temp_dir, f"{uuid.uuid4()}{ext}")
            
            with open(temp_path, "wb") as f:
                f.write(file_bytes)
            
            try:
                text = docx2txt.process(temp_path)
                if not text or not text.strip():
                    return []
                return [{
                    "text": text,
                    "page": 1,
                    "source": filename
                }]
            finally:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
                    
        # 6. PowerPoint (PPTX) Parser (using temp file slide-by-slide)
        elif ext == ".pptx":
            temp_dir = "temp_uploads"
            os.makedirs(temp_dir, exist_ok=True)
            temp_path = os.path.join(temp_dir, f"{uuid.uuid4()}{ext}")
            
            with open(temp_path, "wb") as f:
                f.write(file_bytes)
            
            try:
                prs = Presentation(temp_path)
                slides = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    if slide_text:
                        slides.append({
                            "text": "\n".join(slide_text),
                            "page": i + 1,
                            "source": filename
                        })
                return slides
            finally:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
                    
        else:
            raise ValueError(f"Unsupported file format: {ext}")
            
    except Exception as e:
        logger.error(f"Error parsing file '{filename}': {e}")
        # Raise standard exception to let backend know
        raise e
